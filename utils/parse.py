import os
import io
import requests
import tempfile
import logging
import nest_asyncio
from urllib.parse import urlparse
from pptx import Presentation
from pathlib import Path
from dotenv import load_dotenv

# Tesseract OCR imports
import pytesseract
import httpx
from PIL import Image
from bs4 import BeautifulSoup  # <-- Added for HTML parsing

from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_core.documents import Document as LCDocument
from langchain_community.document_loaders import (
    PyMuPDFLoader,
    Docx2txtLoader,
    UnstructuredEmailLoader,
    CSVLoader,
    UnstructuredExcelLoader
)

# === ENV CONFIG ===
load_dotenv()
nest_asyncio.apply()
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# pytesseract.pytesseract.tesseract_cmd = '/usr/bin/tesseract'
# pytesseract.pytesseract.tesseract_cmd = r'D:\Tesseract\tesseract.exe'
pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"

def extract_text_from_image_with_tesseract(image_path: str) -> str:
    """Extract text from image using Tesseract OCR"""
    try:
        image = Image.open(image_path)
        text = pytesseract.image_to_string(
            image, 
            config='--oem 3 --psm 6'
        )
        logger.info(f"Tesseract OCR successful for {image_path}")
        return text.strip()
    except Exception as e:
        logger.error(f"Tesseract OCR failed for {image_path}: {e}")
        return ""


def extract_images_from_pptx(pptx_path, output_dir="extracted_images"):
    """Extract images from PowerPoint presentation"""
    os.makedirs(output_dir, exist_ok=True)
    prs = Presentation(pptx_path)
    count = 0

    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.shape_type == 13:  # Picture shape type
                image = shape.image
                ext = image.ext
                image_bytes = image.blob
                file_path = os.path.join(output_dir, f"slide_{i+1}_img_{count+1}.{ext}")
                with open(file_path, "wb") as f:
                    f.write(image_bytes)
                count += 1

    return [os.path.join(output_dir, f) for f in os.listdir(output_dir)]


def extract_text_from_pptx_images_with_tesseract(pptx_path: str) -> str:
    """Extract text from all images in a PowerPoint presentation using Tesseract"""
    image_paths = extract_images_from_pptx(pptx_path)
    combined_text = ""

    for image_path in image_paths:
        logger.info(f"Processing image: {image_path}")
        text = extract_text_from_image_with_tesseract(image_path)
        if text:
            combined_text += f"\n--- Image Text ---\n{text}\n"

        if os.path.exists(image_path):
            os.remove(image_path)

    return combined_text


# === Main document parser ===
async def parse_document_from_url(url: str):
    """Parse documents from URL and extract text content (async, single HTTP request)."""
    logger.info(f"Parsing URL: {url}")

    # Single async request
    async with httpx.AsyncClient(follow_redirects=True, timeout=30.0) as client:
        resp = await client.get(url)
        resp.raise_for_status()
        content = resp.content
        content_type = resp.headers.get("Content-Type", "").lower()

    parsed_url = urlparse(url)

    # Determine extension from content type or URL path
    if 'pdf' in content_type:
        ext = '.pdf'
    elif 'wordprocessingml.document' in content_type:
        ext = '.docx'
    elif 'eml' in content_type or 'message/rfc822' in content_type:
        ext = '.eml'
    elif 'jpeg' in content_type or 'jpg' in content_type:
        ext = '.jpg'
    elif 'png' in content_type:
        ext = '.png'
    elif 'pptx' in content_type:
        ext = '.pptx'
    elif 'csv' in content_type:
        ext = '.csv'
    elif 'excel' in content_type or 'spreadsheetml.sheet' in content_type:
        ext = '.xlsx'
    elif 'html' in content_type or 'text/html' in content_type:
        ext = '.html'
    else:
        ext = Path(parsed_url.path).suffix.lower()
        if ext not in [
            '.pdf', '.docx', '.eml', '.csv', '.xlsx', '.pptx',
            '.jpg', '.jpeg', '.png', '.zip', '.html', '.htm'
        ]:
            raise ValueError(f"Unsupported file type: {ext or content_type}")

    # Save to temp file
    with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as tmp_file:
        tmp_file.write(content)
        tmp_path = tmp_file.name

    try:
        # Load based on extension
        if ext == ".pdf":
            loader = PyMuPDFLoader(tmp_path)
            documents = loader.load()

        elif ext == ".docx":
            loader = Docx2txtLoader(tmp_path)
            documents = loader.load()

        elif ext == ".eml":
            loader = UnstructuredEmailLoader(tmp_path)
            documents = loader.load()

        elif ext == ".csv":
            loader = CSVLoader(file_path=tmp_path)
            documents = loader.load()

        elif ext == ".xlsx":
            loader = UnstructuredExcelLoader(file_path=tmp_path)
            documents = loader.load()

        elif ext in ['.jpg', '.jpeg', '.png']:
            text = extract_text_from_image_with_tesseract(tmp_path)
            documents = [LCDocument(page_content=text)]

        elif ext == ".pptx":
            ppt_text = extract_text_from_pptx_images_with_tesseract(tmp_path)
            documents = [LCDocument(page_content=ppt_text)]

        elif ext in ['.html', '.htm']:
            soup = BeautifulSoup(content, "html.parser")
            for tag in soup(["script", "style"]):
                tag.extract()
            text = soup.get_text(separator="\n")
            clean_text = "\n".join(line.strip() for line in text.splitlines() if line.strip())
            documents = [LCDocument(page_content=clean_text)]

        else:
            raise ValueError(f"No loader configured for: {ext}")

        logger.info(documents[0].page_content[:500])
        logger.info(f"Successfully parsed document with {len(documents)} sections")
        return documents

    finally:
        if os.path.exists(tmp_path):
            os.unlink(tmp_path)


# === Split into LangChain-compatible chunks ===
def split_documents(parsed_docs, chunk_size=1000, chunk_overlap=200):
    """Split documents into smaller chunks for processing"""
    all_chunks = []
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap,
        separators=["\n\n", "\n", ".", " "]
    )
    
    try:
        chunks = splitter.split_documents(parsed_docs)
        all_chunks.extend(chunks)
        logger.info(f"Split documents into {len(all_chunks)} chunks")
    except Exception as e:
        logger.error(f"Error processing document chunk: {e}")

    return all_chunks
